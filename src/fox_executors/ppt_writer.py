#!/usr/bin/env python3
"""
PowerPoint 静默生成器 — 接收 JSON 数据，后台生成排版精美的 .pptx 文件。
绝对不弹窗、不抢焦点（不打开 PowerPoint app）。

用法:
    python ppt_writer.py ~/Desktop/报告.pptx '{"title":"...","sections":[...]}'
    python ppt_writer.py ~/Desktop/报告.pptx -  < data.json

JSON 数据结构（与 word_writer 对齐）:
{
  "title":    "2026年4月财务汇总报告",
  "subtitle": "LittleFox 自动生成",
  "sections": [
    {"heading": "一、本月报销概况", "body": "..."},
    {"heading": "二、分类汇总",     "body": "..."}
  ]
}

输出结构:
  Slide 1:  封面（title + subtitle，居中）
  Slide 2-N: 每个 section 一张（heading 顶部，body 占主体）
  Slide N+1: 收尾（"小狐狸 LittleFox · 全场景 AI Agent 桌面操作系统"）
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print(
        "[ppt_writer] ❌ 缺少 python-pptx 库，请运行: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


# 跟 fake invoice 视觉一致的深蓝主色
BRAND_NAVY = RGBColor(0x1C, 0x32, 0x5A)
BRAND_LIGHT_BLUE = RGBColor(0xB4, 0xC8, 0xE6)
TEXT_DARK = RGBColor(0x14, 0x14, 0x14)
TEXT_GREY = RGBColor(0x78, 0x78, 0x78)


def _style_text(run, *, size_pt: int, bold: bool = False, color: RGBColor = TEXT_DARK) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_cover_slide(prs: Presentation, title: str, subtitle: str | None) -> None:
    """封面：深蓝色块 + 大字标题 + 浅色副标题。"""
    blank_layout = prs.slide_layouts[6]  # 6 = Blank
    slide = prs.slides.add_slide(blank_layout)

    # 整页深蓝背景
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_NAVY
    bg.line.fill.background()

    # 大标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        prs.slide_width - Inches(2), Inches(2),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _style_text(run, size_pt=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            prs.slide_width - Inches(2), Inches(1),
        )
        sf = sub_box.text_frame
        sf.word_wrap = True
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = subtitle
        _style_text(srun, size_pt=20, color=BRAND_LIGHT_BLUE)


def _add_content_slide(prs: Presentation, heading: str, body: str) -> None:
    """内容页：顶部深蓝条 + 标题，下面正文。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 顶部色条
    top_band = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(1.2),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = BRAND_NAVY
    top_band.line.fill.background()

    # 标题（在色条里）
    h_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3),
        prs.slide_width - Inches(1.2), Inches(0.8),
    )
    htf = h_box.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hrun = hp.add_run()
    hrun.text = heading
    _style_text(hrun, size_pt=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 正文
    body_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8),
        prs.slide_width - Inches(1.6), prs.slide_height - Inches(2.5),
    )
    btf = body_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    brun = bp.add_run()
    brun.text = body
    _style_text(brun, size_pt=20, color=TEXT_DARK)


def _add_outro_slide(prs: Presentation) -> None:
    """收尾页：品牌签名。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_NAVY
    bg.line.fill.background()

    sig_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.2),
        prs.slide_width - Inches(2), Inches(1.5),
    )
    stf = sig_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = "🦊  Powered by 小狐狸 LittleFox"
    _style_text(srun, size_pt=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    sub_p = stf.add_paragraph()
    sub_p.alignment = PP_ALIGN.CENTER
    sub_run = sub_p.add_run()
    sub_run.text = "全场景 AI Agent 桌面操作系统"
    _style_text(sub_run, size_pt=18, color=BRAND_LIGHT_BLUE)


def run(output_path: str, data_json: str) -> int:
    # ── 解析路径 ──
    target = Path(output_path).expanduser().resolve()
    target.parent.mkdir(parents=True, exist_ok=True)

    # ── 解析 JSON（支持 stdin）──
    if data_json.strip() == "-":
        data_json = sys.stdin.read()
    try:
        data = json.loads(data_json)
    except json.JSONDecodeError as exc:
        print(f"[ppt_writer] ❌ JSON 解析失败: {exc}", file=sys.stderr)
        return 1

    title = data.get("title", "未命名报告").strip() or "未命名报告"
    subtitle = data.get("subtitle")
    sections = data.get("sections") or []
    if not isinstance(sections, list):
        print("[ppt_writer] ❌ 'sections' 必须是数组", file=sys.stderr)
        return 1

    # ── 生成 PPT ──
    try:
        prs = Presentation()
        # 16:9 ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        _add_cover_slide(prs, title, subtitle)
        for sec in sections:
            if not isinstance(sec, dict):
                continue
            heading = str(sec.get("heading") or "").strip()
            body = str(sec.get("body") or "").strip()
            if not heading and not body:
                continue
            _add_content_slide(prs, heading or "(无标题)", body or "(无内容)")
        _add_outro_slide(prs)

        prs.save(str(target))
    except Exception as exc:
        print(f"[ppt_writer] ❌ 生成失败: {exc}", file=sys.stderr)
        return 1

    print(f"[ppt_writer] ✅ 演示文稿已生成 -> {target}")
    return 0


def main() -> int:
    if len(sys.argv) < 3:
        print(
            "用法: python ppt_writer.py <output.pptx> <data_json | ->",
            file=sys.stderr,
        )
        return 1
    return run(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    sys.exit(main())
